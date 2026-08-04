import io
import os
import re
import json
import base64
import tempfile
import sqlite3
import datetime
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# Security & Authentication
from passlib.context import CryptContext
import jwt
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm

from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Depends, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse, JSONResponse
from pydantic import BaseModel

# Gemini AI integration check
try:
    import google.generativeai as genai
    HAS_GEMINI = True
except ImportError:
    HAS_GEMINI = False

# Import PuraLang Core safely
try:
    import puralang
    from puralang import core
    HAS_PURALANG = True
except ImportError:
    HAS_PURALANG = False

# Format-specific optional dependencies
try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pdfplumber
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

# Security & Secret Settings
SECRET_KEY = os.getenv("JWT_SECRET_KEY", "puralang_super_secret_jwt_key_2026")
ALGORITHM = "HS256"
TOKEN_EXPIRE_MINUTES = 60 * 24  # 24 hours

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/api/v1/login", auto_error=False)

# Configure Gemini API Key if present
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
if HAS_GEMINI and GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

app = FastAPI(
    title="PuraLang Cloud API",
    description="Full-stack AI Multi-Format Data & Document Engine",
    version="2.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class PromptRequest(BaseModel):
    prompt: str
    filename: str = "dataset.csv"

class UserRegister(BaseModel):
    username: str
    email: str
    password: str

# Global state for latest processed artifacts
LATEST_CLEANED_DF = None
LATEST_FILENAME = "cleaned_dataset.csv"
LATEST_EXCEL_BYTES = None
LATEST_EXCEL_FILENAME = "processed.xlsx"
LATEST_WORD_BYTES = None
LATEST_WORD_FILENAME = "processed.docx"
LATEST_PPT_BYTES = None
LATEST_PPT_FILENAME = "processed.pptx"
LATEST_PDF_BYTES = None
LATEST_PDF_FILENAME = "processed.pdf"
LATEST_IMAGE_BYTES = None
LATEST_IMAGE_FILENAME = "processed.png"
LATEST_EXTRACTED_DF = None


# ==========================================
# 0. DATABASE SETUP (puralang.db)
# ==========================================

DB_FILE = "puralang.db"

def get_db():
    conn = sqlite3.connect(DB_FILE)
    conn.row_factory = sqlite3.Row
    try:
        yield conn
    finally:
        conn.close()

def init_db():
    conn = sqlite3.connect(DB_FILE)
    cursor = conn.cursor()
    # Users Table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    # Execution & Activity Logs Table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS activity_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            action_type TEXT NOT NULL,
            filename TEXT,
            dsl_rule TEXT,
            timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY(user_id) REFERENCES users(id)
        )
    ''')
    conn.commit()
    conn.close()

init_db()


# ==========================================
# AUTHENTICATION HELPERS
# ==========================================

def hash_password(password: str) -> str:
    return pwd_context.hash(password)

def verify_password(plain_password: str, hashed_password: str) -> bool:
    return pwd_context.verify(plain_password, hashed_password)

def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.datetime.utcnow() + datetime.timedelta(minutes=TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def get_current_user(token: str = Depends(oauth2_scheme), db: sqlite3.Connection = Depends(get_db)):
    if not token:
        return None
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            return None
    except jwt.PyJWTError:
        return None
    
    cursor = db.cursor()
    cursor.execute("SELECT id, username, email, created_at FROM users WHERE username = ?", (username,))
    user = cursor.fetchone()
    return dict(user) if user else None


# ==========================================
# SHARED GEMINI HELPERS
# ==========================================

GEMINI_MODELS = ['gemini-2.5-flash', 'gemini-1.5-flash-latest', 'gemini-1.5-flash', 'gemini-pro']

def call_gemini(system_instruction: str, user_prompt: str) -> str:
    if not HAS_GEMINI or not GEMINI_API_KEY:
        raise HTTPException(status_code=500, detail="Gemini API is not configured.")
    last_error = None
    for model_name in GEMINI_MODELS:
        try:
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(f"{system_instruction}\nUser request: {user_prompt}")
            if response and response.text:
                return response.text.strip()
        except Exception as err:
            last_error = err
            continue
    raise HTTPException(status_code=500, detail=f"Gemini error: {last_error}")

def parse_gemini_json(text: str) -> dict:
    cleaned = re.sub(r"```json|```", "", text).strip()
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", cleaned, re.DOTALL)
        if match:
            return json.loads(match.group())
        return {"raw": text}

def call_gemini_vision(image_bytes: bytes, prompt: str) -> str:
    if not HAS_GEMINI or not GEMINI_API_KEY:
        raise HTTPException(status_code=500, detail="Gemini Vision is not configured.")
    image_part = {"mime_type": "image/png", "data": image_bytes}
    last_error = None
    for model_name in GEMINI_MODELS:
        try:
            model = genai.GenerativeModel(model_name)
            response = model.generate_content([prompt, image_part])
            if response and response.text:
                return response.text.strip()
        except Exception as err:
            last_error = err
            continue
    raise HTTPException(status_code=500, detail=f"Gemini Vision error: {last_error}")


# ==========================================
# 1. AUTH & USER ENDPOINTS
# ==========================================

@app.post("/api/v1/register")
def register_user(user_data: UserRegister, db: sqlite3.Connection = Depends(get_db)):
    cursor = db.cursor()
    cursor.execute("SELECT id FROM users WHERE username = ? OR email = ?", (user_data.username, user_data.email))
    if cursor.fetchone():
        raise HTTPException(status_code=400, detail="Username or Email already registered.")
    
    hashed_pwd = hash_password(user_data.password)
    cursor.execute(
        "INSERT INTO users (username, email, password_hash) VALUES (?, ?, ?)",
        (user_data.username, user_data.email, hashed_pwd)
    )
    db.commit()
    return {"success": True, "message": "Account created successfully!"}

@app.post("/api/v1/login")
def login_for_access_token(form_data: OAuth2PasswordRequestForm = Depends(), db: sqlite3.Connection = Depends(get_db)):
    cursor = db.cursor()
    cursor.execute("SELECT * FROM users WHERE username = ?", (form_data.username,))
    user = cursor.fetchone()
    
    if not user or not verify_password(form_data.password, user["password_hash"]):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    access_token = create_access_token(data={"sub": user["username"]})
    return {"access_token": access_token, "token_type": "bearer", "username": user["username"]}

@app.get("/api/v1/me")
def read_users_me(current_user: dict = Depends(get_current_user)):
    if not current_user:
        raise HTTPException(status_code=401, detail="Not authenticated")
    return current_user

@app.get("/api/v1/admin/logs")
def view_activity_logs(db: sqlite3.Connection = Depends(get_db)):
    cursor = db.cursor()
    cursor.execute("""
        SELECT activity_logs.id, users.username, activity_logs.action_type, 
               activity_logs.filename, activity_logs.dsl_rule, activity_logs.timestamp 
        FROM activity_logs 
        LEFT JOIN users ON activity_logs.user_id = users.id
        ORDER BY activity_logs.timestamp DESC
    """)
    logs = [dict(row) for row in cursor.fetchall()]
    return {"total_logs": len(logs), "logs": logs}


# ==========================================
# 2. FRONTEND ROUTING ENDPOINTS
# ==========================================

@app.get("/")
@app.get("/landing")
@app.get("/landing.html")
def serve_landing():
    if os.path.exists("landing.html"):
        return FileResponse("landing.html")
    elif os.path.exists("index.html"):
        return FileResponse("index.html")
    raise HTTPException(status_code=404, detail="landing.html not found.")

@app.get("/login")
@app.get("/login.html")
def serve_login():
    if os.path.exists("login.html"):
        return FileResponse("login.html")
    elif os.path.exists("landing.html"):
        return FileResponse("landing.html")
    raise HTTPException(status_code=404, detail="login.html not found.")

@app.get("/index")
@app.get("/index.html")
@app.get("/workspace")
def serve_workspace():
    if os.path.exists("index.html"):
        return FileResponse("index.html")
    elif os.path.exists("landing.html"):
        return FileResponse("landing.html")
    raise HTTPException(status_code=404, detail="index.html not found.")

@app.get("/api/v1/health")
def health_check():
    return {
        "status": "online",
        "engine_connected": True,
        "has_gemini": HAS_GEMINI and bool(GEMINI_API_KEY),
        "has_puralang": HAS_PURALANG,
        "formats": {
            "excel": HAS_OPENPYXL,
            "word": HAS_DOCX,
            "ppt": HAS_PPTX,
            "pdf": HAS_PDF,
            "image": HAS_PIL,
            "ocr": HAS_TESSERACT or (HAS_GEMINI and bool(GEMINI_API_KEY)),
        }
    }


# ==========================================
# 3. AI CODE GENERATION
# ==========================================

@app.post("/api/v1/generate")
async def generate_dsl(request: PromptRequest):
    if not request.prompt.strip():
        raise HTTPException(status_code=400, detail="Prompt cannot be empty.")

    target_file = request.filename if request.filename else "dataset.csv"

    system_instruction = (
        "You are an expert compiler for PuraLang DSL.\n"
        "Translate user cleaning requests into valid PuraLang DSL pipelines.\n\n"
        "STRICT PURALANG GRAMMAR RULES:\n"
        f'1. Every script MUST start with: LOAD "{target_file}"\n'
        "2. Steps MUST be chained using the pipe operator: |>\n"
        '3. Column names MUST use double quotes ("...").\n'
        "4. Supported actions:\n"
        '   - DROP_DUPLICATES "column_name"\n'
        '   - FILL_NULLS "column_name" VALUE "default_val"\n'
        '   - FORMAT_STRINGS "column_name" TO UPPERCASE (or LOWERCASE)\n'
        '   - FILTER_ROWS "column_name" == "value"\n'
        '   - RENAME_COLUMN "old_name" TO "new_name"\n'
        '   - SORT_BY "column_name" ASC (or DESC)\n'
        '   - DROP_COLUMN "column_name"\n'
        '   - EXPORT_CSV "output.csv"\n\n'
        "Output ONLY raw DSL code without markdown backticks or explanations."
    )

    cleaned_dsl = call_gemini(system_instruction, request.prompt)
    cleaned_dsl = cleaned_dsl.replace("```puralang", "").replace("```", "").strip()
    return {"success": True, "prompt": request.prompt, "dsl_rule": cleaned_dsl}


# ==========================================
# 4. DSL EXECUTION & ANALYTICS ENDPOINT (CSV)
# ==========================================

@app.post("/api/v1/execute")
async def execute_cleaning(
    file: UploadFile = File(...),
    dsl_rule: str = Form(...),
    db: sqlite3.Connection = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    global LATEST_CLEANED_DF, LATEST_FILENAME

    raw_filename = os.path.basename(file.filename) if file.filename else "temp_dataset.csv"
    temp_filename = f"temp_{raw_filename}"

    try:
        contents = await file.read()

        df_original = None
        try:
            df_original = pd.read_csv(io.BytesIO(contents))
        except Exception:
            try:
                df_original = pd.read_csv(io.BytesIO(contents), encoding="latin1")
            except Exception:
                df_original = pd.read_excel(io.BytesIO(contents))

        header_row_idx = None
        for i in range(min(10, len(df_original))):
            row_str = [str(val).strip() for val in df_original.iloc[i].values]
            if any(col in row_str for col in ["Dept", "Email", "S.No", "ID", "Name"]):
                header_row_idx = i
                break

        if header_row_idx is not None:
            new_cols = [str(val).strip() for val in df_original.iloc[header_row_idx].values]
            df_original = df_original.iloc[header_row_idx + 1:].reset_index(drop=True)
            df_original.columns = new_cols

        df_cleaned = df_original.copy()

        with open(temp_filename, "wb") as f:
            f.write(contents)

        if dsl_rule.strip():
            try:
                filter_matches = re.findall(r'FILTER_ROWS\s+"([^"]+)"\s*(==|!=|>|<)\s*"([^"]+)"', dsl_rule)
                for col_name, op, val in filter_matches:
                    if col_name in df_cleaned.columns:
                        if op == "==":
                            df_cleaned = df_cleaned[df_cleaned[col_name].astype(str).str.strip() == val.strip()]
                        elif op == "!=":
                            df_cleaned = df_cleaned[df_cleaned[col_name].astype(str).str.strip() != val.strip()]

                sort_matches = re.findall(r'SORT_BY\s+"([^"]+)"\s*(ASC|DESC)?', dsl_rule, re.IGNORECASE)
                for col_name, direction in sort_matches:
                    if col_name in df_cleaned.columns:
                        ascending = direction.upper() != "DESC" if direction else True
                        df_cleaned = df_cleaned.sort_values(by=col_name, ascending=ascending)

                drop_matches = re.findall(r'DROP_COLUMN\s+"([^"]+)"', dsl_rule, re.IGNORECASE)
                for col_name in drop_matches:
                    if col_name in df_cleaned.columns:
                        df_cleaned = df_cleaned.drop(columns=[col_name])

                pure_lark_dsl = re.sub(r'\|\>\s*(FILTER_ROWS|SORT_BY|DROP_COLUMN)\s+.*', '', dsl_rule).strip()
                pure_lark_dsl = re.sub(r'LOAD\s+"[^"]+"', f'LOAD "{temp_filename}"', pure_lark_dsl, flags=re.IGNORECASE)

                if HAS_PURALANG and "|>" in pure_lark_dsl and hasattr(core, "PuraTransformer"):
                    parser = core.Lark(core.puralang_grammar, start='start')
                    tree = parser.parse(pure_lark_dsl)
                    transformer = core.PuraTransformer()
                    try:
                        transformed_result = transformer.transform(tree, df_cleaned)
                    except TypeError:
                        transformer.df = df_cleaned
                        transformed_result = transformer.transform(tree)

                    if isinstance(transformed_result, pd.DataFrame):
                        df_cleaned = transformed_result
                    elif hasattr(transformer, "df") and isinstance(transformer.df, pd.DataFrame):
                        df_cleaned = transformer.df

            except Exception as dsl_err:
                raise HTTPException(status_code=400, detail=f"PuraLang Engine Error: {str(dsl_err)}")

        LATEST_CLEANED_DF = df_cleaned.copy()
        LATEST_FILENAME = f"cleaned_{raw_filename}"

        # Log activity to SQLite DB if user is authenticated
        if current_user:
            cursor = db.cursor()
            cursor.execute(
                "INSERT INTO activity_logs (user_id, action_type, filename, dsl_rule) VALUES (?, ?, ?, ?)",
                (current_user["id"], "CSV_DSL_EXECUTION", raw_filename, dsl_rule)
            )
            db.commit()

        chart_base64 = ""
        dept_counts = {}
        target_chart_col = "Dept" if "Dept" in df_original.columns else (df_original.columns[0] if len(df_original.columns) > 0 else None)

        if target_chart_col:
            if target_chart_col in df_cleaned.columns and df_cleaned[target_chart_col].nunique() > 1:
                chart_df = df_cleaned
                title_suffix = " (Filtered Result)"
            else:
                chart_df = df_original
                title_suffix = " (Overall Dataset)"

            if not chart_df.empty:
                counts = chart_df[target_chart_col].value_counts().head(8)
                dept_counts = counts.to_dict()

                fig, ax = plt.subplots(figsize=(6, 3.2), facecolor='#121c17')
                ax.set_facecolor('#121c17')
                colors = ['#00e676', '#10b981', '#3b82f6', '#f59e0b', '#8b5cf6', '#ec4899', '#14b8a6', '#6366f1']
                bars = ax.bar(counts.index.astype(str), counts.values, color=colors[:len(counts)], width=0.45)
                ax.set_title(f'Distribution by {target_chart_col}{title_suffix}', color='#ecf0f1', fontsize=11, fontweight='bold', pad=12)
                ax.tick_params(colors='#95a5a6', labelsize=9)
                plt.xticks(rotation=15, ha='right')
                for bar in bars:
                    yval = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2, yval + (max(counts.values) * 0.02), f"{int(yval)}",
                            ha='center', va='bottom', color='#ecf0f1', fontsize=8, fontweight='bold')
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
                ax.spines['left'].set_color('#1b2e25')
                ax.spines['bottom'].set_color('#1b2e25')
                buf = io.BytesIO()
                plt.tight_layout()
                plt.savefig(buf, format='png', dpi=130)
                plt.close(fig)
                buf.seek(0)
                chart_base64 = base64.b64encode(buf.getvalue()).decode('utf-8')

        return {
            "success": True,
            "filename": raw_filename,
            "metrics": {
                "rows_before": len(df_original),
                "rows_after": len(df_cleaned),
                "rows_dropped": len(df_original) - len(df_cleaned),
                "null_count": int(df_cleaned.isnull().sum().sum())
            },
            "chart_col": target_chart_col,
            "counts": dept_counts,
            "chart_image": chart_base64,
            "dsl_rule": dsl_rule,
            "preview": {
                "before": df_original.fillna("").head(1000).to_dict(orient="records"),
                "after": df_cleaned.fillna("").head(1000).to_dict(orient="records")
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Execution error: {str(e)}")
    finally:
        if os.path.exists(temp_filename):
            os.remove(temp_filename)


# ==========================================
# 5. EXCEL PROCESSING
# ==========================================

def _apply_excel_actions(df: pd.DataFrame, actions: dict) -> pd.DataFrame:
    for rep in actions.get("rename_columns", []):
        old, new = rep.get("from"), rep.get("to")
        if old in df.columns and new:
            df = df.rename(columns={old: new})
    for col in actions.get("drop_columns", []):
        if col in df.columns:
            df = df.drop(columns=[col])
    for filt in actions.get("filter_rows", []):
        col, op, val = filt.get("column"), filt.get("operator", "=="), filt.get("value")
        if col in df.columns:
            if op == "==":
                df = df[df[col].astype(str).str.strip() == str(val).strip()]
            elif op == "contains":
                df = df[df[col].astype(str).str.contains(str(val), case=False, na=False)]
    sort_col = actions.get("sort_by")
    if sort_col and sort_col in df.columns:
        df = df.sort_values(by=sort_col, ascending=actions.get("sort_asc", True))
    return df

@app.post("/api/v1/process-excel")
async def process_excel(
    file: UploadFile = File(...),
    prompt: str = Form(""),
    sheet_name: str = Form(""),
    db: sqlite3.Connection = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    global LATEST_EXCEL_BYTES, LATEST_EXCEL_FILENAME, LATEST_CLEANED_DF

    if not HAS_OPENPYXL:
        raise HTTPException(status_code=500, detail="openpyxl is not installed.")

    contents = await file.read()
    raw_filename = os.path.basename(file.filename) if file.filename else "workbook.xlsx"

    try:
        wb = load_workbook(io.BytesIO(contents), data_only=True)
        sheets = wb.sheetnames
        target_sheet = sheet_name if sheet_name in sheets else sheets[0]

        df = pd.read_excel(io.BytesIO(contents), sheet_name=target_sheet, engine="openpyxl")
        rows_before = len(df)
        df.columns = [str(c).strip() for c in df.columns]
        df = df.dropna(how="all").reset_index(drop=True)

        formulas = []
        ws = wb[target_sheet]
        for row in ws.iter_rows(max_row=min(50, ws.max_row)):
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.startswith("="):
                    formulas.append({"cell": cell.coordinate, "formula": cell.value})
                    if len(formulas) >= 20:
                        break
            if len(formulas) >= 20:
                break

        actions_applied = []
        if prompt.strip() and HAS_GEMINI:
            system = (
                "You translate Excel cleaning requests into JSON actions.\n"
                'Return ONLY JSON: {"rename_columns":[{"from":"A","to":"B"}],'
                '"drop_columns":["col"],"filter_rows":[{"column":"X","operator":"==","value":"Y"}],'
                '"sort_by":"col","sort_asc":true}'
            )
            actions = parse_gemini_json(call_gemini(system, prompt))
            df = _apply_excel_actions(df, actions)
            actions_applied = actions

        out_buf = io.BytesIO()
        with pd.ExcelWriter(out_buf, engine="openpyxl") as writer:
            df.to_excel(writer, sheet_name=target_sheet, index=False)
        out_buf.seek(0)
        LATEST_EXCEL_BYTES = out_buf.getvalue()
        LATEST_EXCEL_FILENAME = f"processed_{raw_filename}"
        LATEST_CLEANED_DF = df.copy()

        if current_user:
            cursor = db.cursor()
            cursor.execute(
                "INSERT INTO activity_logs (user_id, action_type, filename, dsl_rule) VALUES (?, ?, ?, ?)",
                (current_user["id"], "EXCEL_PROCESS", raw_filename, prompt)
            )
            db.commit()

        return {
            "success": True,
            "filename": raw_filename,
            "sheets": sheets,
            "active_sheet": target_sheet,
            "formulas": formulas[:10],
            "actions_applied": actions_applied,
            "metrics": {
                "rows_before": rows_before,
                "rows_after": len(df),
                "columns": len(df.columns),
                "null_count": int(df.isnull().sum().sum())
            },
            "preview": df.fillna("").head(100).to_dict(orient="records")
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Excel processing error: {str(e)}")


# ==========================================
# 6. WORD DOCUMENT PROCESSING
# ==========================================

def _apply_word_actions(doc: "Document", actions: dict) -> int:
    changes = 0
    for rep in actions.get("replace_text", []):
        old, new = rep.get("find", ""), rep.get("replace", "")
        if not old:
            continue
        for para in doc.paragraphs:
            if old in para.text:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        changes += 1
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if old in cell.text:
                        cell.text = cell.text.replace(old, new)
                        changes += 1
    if actions.get("summarize") and HAS_GEMINI:
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if full_text.strip():
            summary = call_gemini("Summarize the document in 3-5 bullet points.", full_text[:8000])
            doc.add_page_break()
            doc.add_heading("AI Summary", level=1)
            for line in summary.split("\n"):
                if line.strip():
                    doc.add_paragraph(line.strip())
            changes += 1
    return changes

@app.post("/api/v1/process-word")
async def process_word(
    file: UploadFile = File(...),
    prompt: str = Form(""),
    db: sqlite3.Connection = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    global LATEST_WORD_BYTES, LATEST_WORD_FILENAME

    if not HAS_DOCX:
        raise HTTPException(status_code=500, detail="python-docx is not installed.")

    contents = await file.read()
    raw_filename = os.path.basename(file.filename) if file.filename else "document.docx"

    try:
        doc = Document(io.BytesIO(contents))
        para_count = len(doc.paragraphs)
        table_count = len(doc.tables)

        actions_applied = {}
        changes = 0
        if prompt.strip() and HAS_GEMINI:
            system = (
                "Translate Word document edit requests into JSON.\n"
                'Return ONLY JSON: {"replace_text":[{"find":"old","replace":"new"}],'
                '"summarize":true/false,"style_headers":true/false}'
            )
            actions_applied = parse_gemini_json(call_gemini(system, prompt))
            changes = _apply_word_actions(doc, actions_applied)

        out_buf = io.BytesIO()
        doc.save(out_buf)
        out_buf.seek(0)
        LATEST_WORD_BYTES = out_buf.getvalue()
        LATEST_WORD_FILENAME = f"processed_{raw_filename}"

        if current_user:
            cursor = db.cursor()
            cursor.execute(
                "INSERT INTO activity_logs (user_id, action_type, filename, dsl_rule) VALUES (?, ?, ?, ?)",
                (current_user["id"], "WORD_PROCESS", raw_filename, prompt)
            )
            db.commit()

        preview_paragraphs = [p.text for p in doc.paragraphs if p.text.strip()][:30]
        preview_tables = []
        for ti, table in enumerate(doc.tables[:3]):
            preview_tables.append([[cell.text for cell in row.cells] for row in table.rows[:5]])

        return {
            "success": True,
            "filename": raw_filename,
            "actions_applied": actions_applied,
            "changes_made": changes,
            "metrics": {
                "paragraphs": para_count,
                "tables": table_count,
                "changes": changes
            },
            "preview": {
                "paragraphs": preview_paragraphs,
                "tables": preview_tables
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Word processing error: {str(e)}")


# ==========================================
# 7. POWERPOINT PROCESSING
# ==========================================

def _apply_ppt_actions(prs: "Presentation", actions: dict) -> int:
    changes = 0
    for rep in actions.get("replace_text", []):
        old, new = rep.get("find", ""), rep.get("replace", "")
        if not old:
            continue
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                changes += 1
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if old in cell.text:
                                cell.text = cell.text.replace(old, new)
                                changes += 1

    if actions.get("add_summary_slide"):
        summary_text = actions.get("summary_text", "AI-generated summary slide.")
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
        tf = txBox.text_frame
        tf.text = "Summary"
        p = tf.add_paragraph()
        p.text = summary_text
        p.font.size = Pt(18)
        changes += 1
    return changes

@app.post("/api/v1/process-ppt")
async def process_ppt(
    file: UploadFile = File(...),
    prompt: str = Form(""),
    db: sqlite3.Connection = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    global LATEST_PPT_BYTES, LATEST_PPT_FILENAME

    if not HAS_PPTX:
        raise HTTPException(status_code=500, detail="python-pptx is not installed.")

    contents = await file.read()
    raw_filename = os.path.basename(file.filename) if file.filename else "presentation.pptx"

    try:
        prs = Presentation(io.BytesIO(contents))
        slide_summaries = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text.strip())
            slide_summaries.append({"slide": i + 1, "text": " | ".join(t for t in texts if t)[:200]})

        actions_applied = {}
        changes = 0
        if prompt.strip() and HAS_GEMINI:
            deck_context = "\n".join(s["text"] for s in slide_summaries[:10])
            system = (
                "Translate PowerPoint edit requests into JSON.\n"
                'Return ONLY JSON: {"replace_text":[{"find":"old","replace":"new"}],'
                '"add_summary_slide":true/false,"summary_text":"text for new slide"}'
            )
            actions_applied = parse_gemini_json(call_gemini(system, f"{prompt}\n\nDeck context:\n{deck_context}"))
            if actions_applied.get("add_summary_slide") and not actions_applied.get("summary_text") and deck_context:
                actions_applied["summary_text"] = call_gemini("Summarize this deck in 4 bullet points.", deck_context)
            changes = _apply_ppt_actions(prs, actions_applied)

        out_buf = io.BytesIO()
        prs.save(out_buf)
        out_buf.seek(0)
        LATEST_PPT_BYTES = out_buf.getvalue()
        LATEST_PPT_FILENAME = f"processed_{raw_filename}"

        if current_user:
            cursor = db.cursor()
            cursor.execute(
                "INSERT INTO activity_logs (user_id, action_type, filename, dsl_rule) VALUES (?, ?, ?, ?)",
                (current_user["id"], "PPT_PROCESS", raw_filename, prompt)
            )
            db.commit()

        return {
            "success": True,
            "filename": raw_filename,
            "slide_count": len(prs.slides),
            "actions_applied": actions_applied,
            "changes_made": changes,
            "metrics": {
                "slides": len(prs.slides),
                "changes": changes
            },
            "preview": {"slides": slide_summaries}
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PowerPoint processing error: {str(e)}")


# ==========================================
# 8. PDF PROCESSING
# ==========================================

def _add_pdf_watermark(input_bytes: bytes, watermark_text: str) -> bytes:
    reader = PdfReader(io.BytesIO(input_bytes))
    writer = PdfWriter()

    watermark_buf = io.BytesIO()
    c = canvas.Canvas(watermark_buf, pagesize=letter)
    c.setFont("Helvetica-Bold", 48)
    c.setFillColorRGB(0.8, 0.1, 0.1, alpha=0.3)
    c.saveState()
    c.translate(300, 400)
    c.rotate(45)
    c.drawCentredString(0, 0, watermark_text)
    c.restoreState()
    c.save()
    watermark_buf.seek(0)
    watermark_reader = PdfReader(watermark_buf)
    watermark_page = watermark_reader.pages[0]

    for page in reader.pages:
        page.merge_page(watermark_page)
        writer.add_page(page)

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return out.getvalue()

@app.post("/api/v1/process-pdf")
async def process_pdf(
    file: UploadFile = File(...),
    prompt: str = Form(""),
    action: str = Form("auto"),
    db: sqlite3.Connection = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    global LATEST_PDF_BYTES, LATEST_PDF_FILENAME, LATEST_EXTRACTED_DF, LATEST_CLEANED_DF

    if not HAS_PDF:
        raise HTTPException(status_code=500, detail="PDF libraries (pypdf, pdfplumber) are not installed.")

    contents = await file.read()
    raw_filename = os.path.basename(file.filename) if file.filename else "document.pdf"

    try:
        reader = PdfReader(io.BytesIO(contents))
        page_count = len(reader.pages)

        extracted_tables = []
        full_text = []
        with pdfplumber.open(io.BytesIO(contents)) as pdf:
            for pi, page in enumerate(pdf.pages[:10]):
                page_text = page.extract_text() or ""
                full_text.append(page_text)
                for table in (page.extract_tables() or [])[:3]:
                    if table and len(table) > 1:
                        headers = [str(h or f"col_{i}") for i, h in enumerate(table[0])]
                        rows = table[1:]
                        extracted_tables.append({
                            "page": pi + 1,
                            "rows": len(rows),
                            "preview": [dict(zip(headers, row)) for row in rows[:5]]
                        })

        actions_applied = {}
        output_bytes = contents

        if prompt.strip() and HAS_GEMINI:
            system = (
                "Translate PDF requests into JSON.\n"
                'Return ONLY JSON: {"action":"extract_tables"|"watermark"|"replace_text",'
                '"watermark_text":"CONFIDENTIAL","replace_text":[{"find":"x","replace":"y"}]}'
            )
            actions_applied = parse_gemini_json(call_gemini(system, prompt))
            action = actions_applied.get("action", action)

        if action in ("watermark", "auto") and (actions_applied.get("watermark_text") or "watermark" in prompt.lower() or "confidential" in prompt.lower() or "approved" in prompt.lower()):
            wm = actions_applied.get("watermark_text", "")
            if not wm:
                wm = "CONFIDENTIAL" if "confidential" in prompt.lower() else "APPROVED" if "approved" in prompt.lower() else "CONFIDENTIAL"
            output_bytes = _add_pdf_watermark(contents, wm)
            actions_applied["watermark_applied"] = wm

        if extracted_tables:
            first = extracted_tables[0]["preview"]
            if first:
                LATEST_EXTRACTED_DF = pd.DataFrame(first)
                LATEST_CLEANED_DF = LATEST_EXTRACTED_DF.copy()

        LATEST_PDF_BYTES = output_bytes
        LATEST_PDF_FILENAME = f"processed_{raw_filename}"

        if current_user:
            cursor = db.cursor()
            cursor.execute(
                "INSERT INTO activity_logs (user_id, action_type, filename, dsl_rule) VALUES (?, ?, ?, ?)",
                (current_user["id"], "PDF_PROCESS", raw_filename, prompt)
            )
            db.commit()

        return {
            "success": True,
            "filename": raw_filename,
            "actions_applied": actions_applied,
            "metrics": {
                "pages": page_count,
                "tables_found": len(extracted_tables),
                "text_length": sum(len(t) for t in full_text)
            },
            "preview": {
                "text_snippet": "\n".join(full_text)[:2000],
                "tables": extracted_tables
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF processing error: {str(e)}")


# ==========================================
# 9. IMAGE / OCR PROCESSING
# ==========================================

@app.post("/api/v1/process-image")
async def process_image(
    file: UploadFile = File(...),
    prompt: str = Form(""),
    db: sqlite3.Connection = Depends(get_db),
    current_user: dict = Depends(get_current_user)
):
    global LATEST_IMAGE_BYTES, LATEST_IMAGE_FILENAME, LATEST_EXTRACTED_DF, LATEST_CLEANED_DF

    if not HAS_PIL:
        raise HTTPException(status_code=500, detail="Pillow is not installed.")

    contents = await file.read()
    raw_filename = os.path.basename(file.filename) if file.filename else "image.png"

    try:
        img = Image.open(io.BytesIO(contents)).convert("RGB")
        width, height = img.size

        ocr_text = ""
        if HAS_TESSERACT:
            try:
                ocr_text = pytesseract.image_to_string(img)
            except Exception:
                ocr_text = ""

        ai_analysis = ""
        if HAS_GEMINI and GEMINI_API_KEY:
            img_buf = io.BytesIO()
            img.save(img_buf, format="PNG")
            img_buf.seek(0)
            vision_prompt = prompt.strip() or "Extract all visible text and describe any tables. Return structured plain text."
            ai_analysis = call_gemini_vision(img_buf.getvalue(), vision_prompt)

        combined_text = ai_analysis or ocr_text

        annotated = img.copy()
        draw = ImageDraw.Draw(annotated)
        draw.rectangle([10, 10, min(width - 10, 320), 50], fill=(18, 28, 23))
        draw.text((20, 20), "PuraLang OCR", fill=(0, 230, 118))

        out_buf = io.BytesIO()
        annotated.save(out_buf, format="PNG")
        out_buf.seek(0)
        LATEST_IMAGE_BYTES = out_buf.getvalue()
        LATEST_IMAGE_FILENAME = f"annotated_{os.path.splitext(raw_filename)[0]}.png"

        table_rows = []
        if combined_text and HAS_GEMINI:
            table_json = parse_gemini_json(call_gemini(
                'If the text contains tabular data, return JSON {"rows":[{"col1":"val",...},...]}. Else {"rows":[]}.',
                combined_text[:6000]
            ))
            table_rows = table_json.get("rows", [])
            if table_rows:
                LATEST_EXTRACTED_DF = pd.DataFrame(table_rows)
                LATEST_CLEANED_DF = LATEST_EXTRACTED_DF.copy()

        preview_b64 = base64.b64encode(out_buf.getvalue()).decode("utf-8")

        if current_user:
            cursor = db.cursor()
            cursor.execute(
                "INSERT INTO activity_logs (user_id, action_type, filename, dsl_rule) VALUES (?, ?, ?, ?)",
                (current_user["id"], "IMAGE_OCR_PROCESS", raw_filename, prompt)
            )
            db.commit()

        return {
            "success": True,
            "filename": raw_filename,
            "metrics": {
                "width": width,
                "height": height,
                "ocr_chars": len(ocr_text),
                "ai_chars": len(ai_analysis),
                "table_rows": len(table_rows)
            },
            "preview": {
                "ocr_text": ocr_text[:3000],
                "ai_analysis": ai_analysis[:3000],
                "table": table_rows[:20],
                "annotated_image": preview_b64
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Image processing error: {str(e)}")


# ==========================================
# 10. DOWNLOAD ENDPOINTS
# ==========================================

@app.get("/api/v1/download")
def download_cleaned_csv():
    global LATEST_CLEANED_DF, LATEST_FILENAME
    if LATEST_CLEANED_DF is None:
        raise HTTPException(status_code=400, detail="No cleaned dataset available. Run a pipeline first.")
    stream = io.StringIO()
    LATEST_CLEANED_DF.to_csv(stream, index=False)
    response = StreamingResponse(iter([stream.getvalue()]), media_type="text/csv")
    response.headers["Content-Disposition"] = f"attachment; filename={LATEST_FILENAME}"
    return response

@app.get("/api/v1/download-excel")
def download_excel():
    global LATEST_EXCEL_BYTES, LATEST_EXCEL_FILENAME
    if not LATEST_EXCEL_BYTES:
        raise HTTPException(status_code=400, detail="No processed Excel file available.")
    return StreamingResponse(
        io.BytesIO(LATEST_EXCEL_BYTES),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={LATEST_EXCEL_FILENAME}"}
    )

@app.get("/api/v1/download-word")
def download_word():
    global LATEST_WORD_BYTES, LATEST_WORD_FILENAME
    if not LATEST_WORD_BYTES:
        raise HTTPException(status_code=400, detail="No processed Word file available.")
    return StreamingResponse(
        io.BytesIO(LATEST_WORD_BYTES),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={LATEST_WORD_FILENAME}"}
    )

@app.get("/api/v1/download-ppt")
def download_ppt():
    global LATEST_PPT_BYTES, LATEST_PPT_FILENAME
    if not LATEST_PPT_BYTES:
        raise HTTPException(status_code=400, detail="No processed PowerPoint file available.")
    return StreamingResponse(
        io.BytesIO(LATEST_PPT_BYTES),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={LATEST_PPT_FILENAME}"}
    )

@app.get("/api/v1/download-pdf")
def download_pdf():
    global LATEST_PDF_BYTES, LATEST_PDF_FILENAME
    if not LATEST_PDF_BYTES:
        raise HTTPException(status_code=400, detail="No processed PDF file available.")
    return StreamingResponse(
        io.BytesIO(LATEST_PDF_BYTES),
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={LATEST_PDF_FILENAME}"}
    )

@app.get("/api/v1/download-image")
def download_image():
    global LATEST_IMAGE_BYTES, LATEST_IMAGE_FILENAME
    if not LATEST_IMAGE_BYTES:
        raise HTTPException(status_code=400, detail="No processed image available.")
    return StreamingResponse(
        io.BytesIO(LATEST_IMAGE_BYTES),
        media_type="image/png",
        headers={"Content-Disposition": f"attachment; filename={LATEST_IMAGE_FILENAME}"}
    )

@app.get("/api/v1/download-extracted")
def download_extracted_table():
    global LATEST_EXTRACTED_DF
    if LATEST_EXTRACTED_DF is None:
        raise HTTPException(status_code=400, detail="No extracted table available.")
    stream = io.StringIO()
    LATEST_EXTRACTED_DF.to_csv(stream, index=False)
    response = StreamingResponse(iter([stream.getvalue()]), media_type="text/csv")
    response.headers["Content-Disposition"] = "attachment; filename=extracted_table.csv"
    return response


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)